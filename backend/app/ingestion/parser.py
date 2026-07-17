"""Document parsing: PDF, DOCX, PPTX, TXT, Markdown, Images."""
from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Optional

from app.core.logging import get_logger

logger = get_logger(__name__)


class ParsedDocument:
    def __init__(self, text: str, pages: list[dict], metadata: dict):
        self.text = text
        self.pages = pages          # [{page_num, text, tables}]
        self.metadata = metadata    # author, title, etc.


def _parse_pdf(file_path: Path) -> ParsedDocument:
    try:
        import pdfplumber
        pages = []
        full_text = []
        meta = {}
        with pdfplumber.open(str(file_path)) as pdf:
            meta = pdf.metadata or {}
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                tables = page.extract_tables() or []
                table_text = _tables_to_text(tables)
                combined = text + ("\n\n" + table_text if table_text else "")
                pages.append({"page_num": i, "text": combined, "tables": tables})
                full_text.append(combined)
        return ParsedDocument(
            text="\n\n".join(full_text),
            pages=pages,
            metadata={
                "title": meta.get("Title", ""),
                "author": meta.get("Author", ""),
                "total_pages": len(pages),
            },
        )
    except Exception as e:
        logger.error("PDF parse failed", path=str(file_path), error=str(e))
        raise


def _parse_docx(file_path: Path) -> ParsedDocument:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    text = "\n\n".join(paragraphs)
    return ParsedDocument(
        text=text,
        pages=[{"page_num": 1, "text": text, "tables": []}],
        metadata={"title": doc.core_properties.title or "", "author": doc.core_properties.author or ""},
    )


def _parse_pptx(file_path: Path) -> ParsedDocument:
    from pptx import Presentation
    prs = Presentation(str(file_path))
    pages = []
    full_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        combined = "\n".join(slide_text)
        pages.append({"page_num": i, "text": combined, "tables": []})
        full_text.append(f"[Slide {i}]\n{combined}")
    return ParsedDocument(text="\n\n".join(full_text), pages=pages, metadata={"total_pages": len(pages)})


def _parse_text(file_path: Path) -> ParsedDocument:
    import chardet
    raw = file_path.read_bytes()
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    text = raw.decode(encoding, errors="replace")
    return ParsedDocument(
        text=text,
        pages=[{"page_num": 1, "text": text, "tables": []}],
        metadata={},
    )


def _parse_image(file_path: Path) -> ParsedDocument:
    """OCR using Tesseract."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(str(file_path))
        text = pytesseract.image_to_string(img)
        return ParsedDocument(
            text=text,
            pages=[{"page_num": 1, "text": text, "tables": []}],
            metadata={"ocr": True},
        )
    except Exception as e:
        logger.warning("OCR failed, returning empty", error=str(e))
        return ParsedDocument(text="", pages=[], metadata={"ocr_error": str(e)})


def _tables_to_text(tables: list) -> str:
    lines = []
    for table in tables:
        for row in table:
            row_cells = [str(c or "").strip() for c in row]
            lines.append(" | ".join(row_cells))
        lines.append("")
    return "\n".join(lines)


def parse_document(file_path: str | Path) -> ParsedDocument:
    path = Path(file_path)
    ext = path.suffix.lower()
    logger.info("Parsing document", path=str(path), ext=ext)

    parsers = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".txt": _parse_text,
        ".md": _parse_text,
        ".markdown": _parse_text,
        ".png": _parse_image,
        ".jpg": _parse_image,
        ".jpeg": _parse_image,
        ".tiff": _parse_image,
        ".bmp": _parse_image,
    }
    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return parser(path)
